import os
from datetime import datetime
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import cm
from pptx import Presentation
from pptx.util import Inches
def make_pdf_report(metrics: dict, path):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    c = canvas.Canvas(path, pagesize=A4); w,h = A4
    c.setFont("Helvetica-Bold",16); c.drawString(2*cm,h-2*cm,"Water Quality Classifier – Summary")
    c.setFont("Helvetica",11); y=h-3*cm
    c.drawString(2*cm,y,f"Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"); y-=1*cm
    for k,v in metrics.items(): c.drawString(2*cm,y,f"{k}: {v}"); y-=0.7*cm
    c.showPage(); c.save(); return path
def make_pptx_deck(metrics: dict, path):
    from pptx.util import Pt
    prs = Presentation(); s=prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text="Water Quality Classifier – Results"; s.placeholders[1].text="Auto-generated summary"
    s2=prs.slides.add_slide(prs.slide_layouts[5]); tx=s2.shapes.add_textbox(Inches(1),Inches(1),Inches(8),Inches(5)).text_frame
    tx.text="Metrics"; 
    for k,v in metrics.items(): p=tx.add_paragraph(); p.text=f"{k}: {v}"; p.level=1
    prs.save(path); return path
